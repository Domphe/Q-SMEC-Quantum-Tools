import os
import sys
import json
import csv
import hashlib
from datetime import datetime

# Optional imports for rich parsing
DOCX_ENABLED = False
PPTX_ENABLED = False
try:
    import docx  # python-docx
    DOCX_ENABLED = True
except Exception:
    pass
try:
    from pptx import Presentation  # python-pptx
    PPTX_ENABLED = True
except Exception:
    pass

ROOTS = [
    r"G:\\My Drive\\White Paper (Patent)",
    r"G:\\My Drive\\NDA Documents",
    r"G:\\My Drive\\Pitch Decks",
    r"G:\\My Drive\\Q-SMEC Overall DB",
]

OUTPUT_DIR = r"G:\\My Drive\\Databases\\QCBD\\reports"
TS = datetime.now().strftime("%Y%m%d_%H%M%S")

SUMMARY_JSON = os.path.join(OUTPUT_DIR, f"deep_summary_{TS}.json")
INVENTORY_CSV = os.path.join(OUTPUT_DIR, f"deep_inventory_{TS}.csv")
TREE_JSON = os.path.join(OUTPUT_DIR, f"deep_trees_{TS}.json")

TEXT_EXTS = {".txt", ".md", ".json", ".csv"}
DOCX_EXTS = {".docx"}
PPTX_EXTS = {".pptx"}


def sha1_file(path: str) -> str:
    try:
        h = hashlib.sha1()
        with open(path, "rb") as f:
            for chunk in iter(lambda: f.read(8192), b""):
                h.update(chunk)
        return h.hexdigest()
    except Exception:
        return ""


def read_text_file(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()
    except Exception:
        return ""


def read_docx(path: str) -> str:
    if not DOCX_ENABLED:
        return ""
    try:
        doc = docx.Document(path)
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception:
        return ""


def read_pptx(path: str) -> str:
    if not PPTX_ENABLED:
        return ""
    try:
        prs = Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception:
        return ""


def summarize_text(text: str) -> dict:
    # Very lightweight heuristic summary; full NLP can be added later
    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
    sample = lines[:25]
    keys = []
    for kw in ["GaN", "AlN", "Ge", "Nb", "Ti", "graphene", "battery", "THz", "defect", "superconduct", "DFT", "HSE06", "QMC", "DOE"]:
        if kw.lower() in text.lower():
            keys.append(kw)
    return {
        "line_count": len(lines),
        "sample_head": sample,
        "keywords": sorted(set(keys)),
    }


def walk_root(root: str):
    tree = {"path": root, "children": []}
    inventory_rows = []
    summaries = []
    for dirpath, dirnames, filenames in os.walk(root):
        # Tree structure
        rel_children = [{"dir": os.path.join(dirpath, d)} for d in dirnames]
        if rel_children:
            tree["children"].extend(rel_children)
        for fn in filenames:
            path = os.path.join(dirpath, fn)
            ext = os.path.splitext(fn)[1].lower()
            try:
                stat = os.stat(path)
                size = stat.st_size
                mtime = datetime.fromtimestamp(stat.st_mtime).isoformat()
            except Exception:
                size = None
                mtime = None
            sha1 = sha1_file(path)
            inventory_rows.append({
                "root": root,
                "path": path,
                "ext": ext,
                "size": size,
                "mtime": mtime,
                "sha1": sha1,
            })
            # Read and summarize content when feasible
            text = ""
            if ext in TEXT_EXTS:
                text = read_text_file(path)
            elif ext in DOCX_EXTS:
                text = read_docx(path)
            elif ext in PPTX_EXTS:
                text = read_pptx(path)
            if text:
                summaries.append({
                    "path": path,
                    "ext": ext,
                    "summary": summarize_text(text),
                })
    return tree, inventory_rows, summaries


def main():
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    all_trees = {}
    all_inventory = []
    all_summaries = []

    for root in ROOTS:
        tree, inv, sums = walk_root(root)
        all_trees[root] = tree
        all_inventory.extend(inv)
        all_summaries.extend(sums)
        print(f"[OK] Crawled {root}: {len(inv)} files, {len(sums)} summarized")

    # Write CSV inventory
    with open(INVENTORY_CSV, "w", newline="", encoding="utf-8") as f:
        w = csv.DictWriter(f, fieldnames=["root", "path", "ext", "size", "mtime", "sha1"])
        w.writeheader()
        for row in all_inventory:
            w.writerow(row)
    print(f"[WRITE] Inventory: {INVENTORY_CSV} ({len(all_inventory)} rows)")

    # Write JSON trees and summaries
    with open(TREE_JSON, "w", encoding="utf-8") as f:
        json.dump(all_trees, f, indent=2)
    print(f"[WRITE] Trees: {TREE_JSON}")

    with open(SUMMARY_JSON, "w", encoding="utf-8") as f:
        json.dump({"generated": TS, "summaries": all_summaries}, f, indent=2)
    print(f"[WRITE] Summaries: {SUMMARY_JSON} ({len(all_summaries)} docs)")

    # Guidance for missing optional parsers
    if not DOCX_ENABLED:
        print("[INFO] python-docx not installed; DOCX content parsed as metadata only.")
    if not PPTX_ENABLED:
        print("[INFO] python-pptx not installed; PPTX content parsed as metadata only.")

    return 0

if __name__ == "__main__":
    sys.exit(main())
